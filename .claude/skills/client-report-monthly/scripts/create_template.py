"""
create_template.py — Tạo file monthly_report_template.pptx lần đầu.

Chạy một lần để sinh template:
    python scripts/create_template.py
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SKILL_DIR = Path(__file__).resolve().parent.parent
TEMPLATE_PATH = SKILL_DIR / "templates" / "monthly_report_template.pptx"
COLORS_PATH = SKILL_DIR / "templates" / "brand_assets" / "colors.json"


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_text_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    bold: bool = False,
    color: str = "#1A1A2E",
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Thêm text box vào slide."""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)


def add_colored_rect(slide, left, top, width, height, fill_color: str) -> None:
    """Thêm hình chữ nhật màu."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.color.rgb = hex_to_rgb(fill_color)


def main() -> None:
    colors = json.loads(COLORS_PATH.read_text())
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ------------------------------------------------------------------
    # Slide 1: Cover
    # ------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_colored_rect(slide1, 0, 0, 13.33, 7.5, colors["primary"])
    add_colored_rect(slide1, 0, 5.8, 13.33, 1.7, colors["accent"])
    add_text_box(slide1, "SEO PERFORMANCE REPORT", 1, 1.5, 11, 1, 32, True, "#FFFFFF", PP_ALIGN.CENTER)
    add_text_box(slide1, "{{CLIENT_NAME}}", 1, 2.8, 11, 1, 40, True, "#FFFFFF", PP_ALIGN.CENTER)
    add_text_box(slide1, "{{PERIOD}}", 1, 4.2, 11, 0.6, 20, False, "#FFFFFF", PP_ALIGN.CENTER)
    add_text_box(slide1, "Prepared by SEONGON Agency", 1, 6.0, 11, 0.5, 14, False, "#FFFFFF", PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # Slide 2: Executive Summary
    # ------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_colored_rect(slide2, 0, 0, 13.33, 1.0, colors["primary"])
    add_text_box(slide2, "EXECUTIVE SUMMARY", 0.3, 0.1, 12, 0.8, 24, True, "#FFFFFF")
    add_text_box(slide2, "{{SUMMARY_HIGHLIGHTS}}", 0.5, 1.2, 12, 5.5, 16, False, colors["text_dark"])
    # KPI boxes
    for i, (label, placeholder) in enumerate([
        ("Organic Sessions", "{{GA4_SESSIONS}}"),
        ("GSC Clicks", "{{GSC_CLICKS}}"),
        ("Conversions", "{{GA4_CONVERSIONS}}"),
        ("Domain Rating", "{{AHREFS_DR}}"),
    ]):
        x = 0.3 + i * 3.2
        add_colored_rect(slide2, x, 5.2, 2.9, 1.8, colors["background_alt"])
        add_text_box(slide2, label, x + 0.1, 5.3, 2.7, 0.5, 12, False, colors["primary"])
        add_text_box(slide2, placeholder, x + 0.1, 5.7, 2.7, 0.8, 28, True, colors["accent"])

    # ------------------------------------------------------------------
    # Slide 3: GA4 Organic Traffic
    # ------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_colored_rect(slide3, 0, 0, 13.33, 1.0, colors["secondary"])
    add_text_box(slide3, "ORGANIC TRAFFIC OVERVIEW (GA4)", 0.3, 0.1, 12, 0.8, 22, True, "#FFFFFF")
    add_text_box(slide3, "Sessions: {{GA4_SESSIONS}} | MoM: {{GA4_SESSIONS_MOM}}", 0.5, 1.2, 12, 0.6, 18, True, colors["primary"])
    add_text_box(slide3, "{{GA4_TOP_PAGES_TABLE}}", 0.5, 2.0, 12, 4.5, 14, False, colors["text_dark"])

    # ------------------------------------------------------------------
    # Slide 4: GSC Performance
    # ------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_colored_rect(slide4, 0, 0, 13.33, 1.0, colors["primary"])
    add_text_box(slide4, "GOOGLE SEARCH CONSOLE PERFORMANCE", 0.3, 0.1, 12, 0.8, 22, True, "#FFFFFF")
    for i, (label, placeholder) in enumerate([
        ("Total Clicks", "{{GSC_CLICKS}}"),
        ("Impressions", "{{GSC_IMPRESSIONS}}"),
        ("Avg CTR", "{{GSC_CTR}}"),
        ("Avg Position", "{{GSC_POSITION}}"),
    ]):
        x = 0.3 + i * 3.2
        add_colored_rect(slide4, x, 1.2, 2.9, 1.5, colors["background_alt"])
        add_text_box(slide4, label, x + 0.1, 1.3, 2.7, 0.5, 12, False, colors["primary"])
        add_text_box(slide4, placeholder, x + 0.1, 1.7, 2.7, 0.7, 26, True, colors["accent"])
    add_text_box(slide4, "Top Queries:\n{{GSC_TOP_QUERIES_TABLE}}", 0.5, 3.0, 12, 4.0, 14, False, colors["text_dark"])

    # ------------------------------------------------------------------
    # Slide 5: Top Performing Pages
    # ------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_colored_rect(slide5, 0, 0, 13.33, 1.0, colors["secondary"])
    add_text_box(slide5, "TOP PERFORMING PAGES", 0.3, 0.1, 12, 0.8, 22, True, "#FFFFFF")
    add_text_box(slide5, "{{TOP_PAGES_TABLE}}", 0.5, 1.2, 12, 5.8, 13, False, colors["text_dark"])

    # ------------------------------------------------------------------
    # Slide 6: Backlinks Growth
    # ------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_colored_rect(slide6, 0, 0, 13.33, 1.0, colors["primary"])
    add_text_box(slide6, "BACKLINKS & AUTHORITY GROWTH (AHREFS)", 0.3, 0.1, 12, 0.8, 22, True, "#FFFFFF")
    for i, (label, placeholder) in enumerate([
        ("Domain Rating", "{{AHREFS_DR}}"),
        ("Referring Domains", "{{AHREFS_REF_DOMAINS}}"),
        ("New Backlinks", "{{AHREFS_NEW_BL}}"),
    ]):
        x = 0.5 + i * 4.2
        add_colored_rect(slide6, x, 1.2, 3.8, 1.5, colors["background_alt"])
        add_text_box(slide6, label, x + 0.1, 1.3, 3.6, 0.5, 12, False, colors["primary"])
        add_text_box(slide6, placeholder, x + 0.1, 1.7, 3.6, 0.7, 28, True, colors["accent"])
    add_text_box(slide6, "New Notable Backlinks:\n{{NEW_BACKLINKS_LIST}}", 0.5, 3.0, 12, 4.0, 14, False, colors["text_dark"])

    # ------------------------------------------------------------------
    # Slide 7: Next Month Action Plan
    # ------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_colored_rect(slide7, 0, 0, 13.33, 1.0, colors["accent"])
    add_text_box(slide7, "NEXT MONTH ACTION PLAN", 0.3, 0.1, 12, 0.8, 22, True, "#FFFFFF")
    add_text_box(slide7, "{{ACTION_PLAN_LIST}}", 0.5, 1.2, 12, 5.5, 16, False, colors["text_dark"])

    # ------------------------------------------------------------------
    prs.save(TEMPLATE_PATH)
    print(f"✅ Template created: {TEMPLATE_PATH}")


if __name__ == "__main__":
    main()
