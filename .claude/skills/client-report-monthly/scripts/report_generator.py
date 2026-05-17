"""
report_generator.py — Tạo monthly SEO report .pptx 15 slides.

Usage (mock):
    python scripts/report_generator.py \\
        --client "Viettel Store" --domain viettelstore.vn \\
        --month 2026-04 --brand viettel_store --mock
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SKILL_DIR = Path(__file__).resolve().parent.parent
WORKSPACE_ROOT = SKILL_DIR.parent.parent.parent
sys.path.insert(0, str(SKILL_DIR / "scripts"))

from brand_loader import (
    load_brand, get_color, get_font, get_logo_path,
    format_vn_number, format_vn_pct, format_delta, lighten_hex,
)

MOCK_DIR    = SKILL_DIR / "tests" / "mock_data"
OUTPUTS_DIR = WORKSPACE_ROOT / "outputs"

# ---------------------------------------------------------------------------
# Kích thước slide 16:9
# ---------------------------------------------------------------------------
W      = Inches(13.333)
H      = Inches(7.5)
MARGIN = Inches(0.5)

# Layout constants dùng chung cho 4-zone slides (6, 8, 9, 10)
_ZLW  = Inches(6.067)                        # zone width (mỗi cột)
_ZRX  = MARGIN + Inches(6.067) + Inches(0.2) # right zone x-start
_ZTY  = Inches(1.15)                         # top zones y-start
_ZTH  = Inches(2.65)                         # top zones height
_ZBY  = Inches(3.9)                          # bottom zones y-start
_ZBH  = Inches(2.55)                         # bottom zones height
_ZNY  = Inches(6.55)                         # narrative y
_ZNH  = Inches(0.42)                         # narrative height

# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _lighten(hex_color: str, factor: float = 0.12) -> str:
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return "#{:02X}{:02X}{:02X}".format(
        min(255, round(r + (255 - r) * factor)),
        min(255, round(g + (255 - g) * factor)),
        min(255, round(b + (255 - b) * factor)),
    )


def _add_rect(slide, left, top, width, height, fill_hex: str, show_line: bool = False):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if show_line:
        shape.line.color.rgb = _rgb(fill_hex)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, left, top, width, height, text, *,
              font_name="Arial", size=14, bold=False,
              color_hex="#222222", align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.name      = font_name
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = _rgb(color_hex)
    return box


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _footer(slide, brand, label: str = ""):
    font_b     = get_font(brand, "body")
    text_light = get_color(brand, "text_light")
    _add_rect(slide, 0, H - Inches(0.36), W, Pt(1), text_light)
    txt = label or f"{brand['client_name']}  •  Báo cáo SEO  •  seongon.com"
    _add_text(slide, MARGIN, H - Inches(0.34), W - 2*MARGIN, Inches(0.32),
              txt, font_name=font_b, size=9,
              color_hex=text_light, align=PP_ALIGN.CENTER)


def _format_month_vn(year_month: str) -> str:
    y, m = year_month.split("-")
    return f"Tháng {int(m)}/{y}"


# ---------------------------------------------------------------------------
# Slide-level helpers (dùng chung cho nhiều slides)
# ---------------------------------------------------------------------------

def _slide_base(slide, brand, title: str):
    """Nền trắng + thanh accent trái + tiêu đề + gạch dưới."""
    primary = get_color(brand, "primary")
    accent  = get_color(brand, "accent")
    bg      = get_color(brand, "background")
    font_h  = get_font(brand, "heading")
    _add_rect(slide, 0, 0, W, H, bg)
    _add_rect(slide, 0, 0, Inches(0.2), H, primary)
    _add_text(slide, MARGIN, Inches(0.32), W - 2*MARGIN, Inches(0.72),
              title, font_name=font_h, size=26, bold=True, color_hex=primary)
    _add_rect(slide, MARGIN, Inches(1.04), Inches(5), Pt(3), accent)


def _zone_hdr(slide, x, y, w, icon: str, title: str, bg_hex: str, brand):
    """Header strip cho zone (màu đặc + icon + text trắng)."""
    font_h = get_font(brand, "heading")
    _add_rect(slide, x, y, w, Inches(0.4), bg_hex)
    _add_text(slide, x + Inches(0.12), y + Inches(0.04),
              w - Inches(0.2), Inches(0.34),
              f"{icon}  {title}", font_name=font_h, size=12,
              bold=True, color_hex="#FFFFFF")


def _narr(slide, brand, text: str):
    """Một dòng narrative ở cuối slide (trước footer)."""
    font_b    = get_font(brand, "body")
    text_dark = get_color(brand, "text_dark")
    _add_rect(slide, MARGIN, _ZNY, W - 2*MARGIN, _ZNH, "#EEF2F7")
    _add_text(slide, MARGIN + Inches(0.12), _ZNY + Inches(0.04),
              W - 2*MARGIN - Inches(0.2), _ZNH - Inches(0.08),
              f"💡  {text}", font_name=font_b, size=10,
              color_hex=text_dark, wrap=True)


def _calc_delta(cur, prev, brand, higher_is_better: bool = True,
                pct: bool = True) -> tuple[str, str]:
    """(delta_str, hex_color). Ví dụ: ('▲ +15,1%', '#28A745')"""
    try:
        c, p = float(cur), float(prev)
        if p == 0:
            return "N/A", get_color(brand, "text_light")
        if pct:
            delta  = (c - p) / p
            prefix = "▲ +" if delta >= 0 else "▼ "
            txt    = f"{prefix}{abs(delta)*100:.1f}%".replace(".", ",")
        else:
            diff   = c - p
            prefix = "▲ +" if diff >= 0 else "▼ "
            txt    = f"{prefix}{abs(diff):.1f}".replace(".", ",")
        is_good = (delta >= 0 if pct else diff >= 0) if higher_is_better \
                  else (delta <= 0 if pct else diff <= 0)
        color = get_color(brand, "success") if is_good else get_color(brand, "danger")
        return txt, color
    except Exception:
        return "N/A", get_color(brand, "text_light")


def _kpi_card(slide, x, y, w, h, label: str, value_str: str,
              delta_str: str, delta_hex: str, brand, *, top_bar_hex: str = None):
    """KPI card: label | big number | delta | 'So với tháng trước'."""
    font_h     = get_font(brand, "heading")
    font_b     = get_font(brand, "body")
    text_dark  = get_color(brand, "text_dark")
    text_light = get_color(brand, "text_light")
    bg_hex     = lighten_hex(get_color(brand, "primary"), 0.92)

    _add_rect(slide, x, y, w, h, bg_hex)

    # Top color stripe
    bar_color = top_bar_hex or (get_color(brand, "success")
                                if "▲" in delta_str else get_color(brand, "danger"))
    _add_rect(slide, x, y, w, Pt(5), bar_color)

    # Label
    _add_text(slide, x + Inches(0.1), y + Inches(0.18),
              w - Inches(0.2), Inches(0.42),
              label, font_name=font_b, size=11,
              color_hex=text_light, align=PP_ALIGN.CENTER)

    # Big number
    _add_text(slide, x + Inches(0.06), y + Inches(0.58),
              w - Inches(0.12), Inches(0.85),
              value_str, font_name=font_h, size=30, bold=True,
              color_hex=text_dark, align=PP_ALIGN.CENTER)

    # Delta
    _add_text(slide, x + Inches(0.06), y + h - Inches(0.72),
              w - Inches(0.12), Inches(0.38),
              delta_str, font_name=font_b, size=13, bold=True,
              color_hex=delta_hex, align=PP_ALIGN.CENTER)

    # "So với tháng trước"
    _add_text(slide, x + Inches(0.06), y + h - Inches(0.38),
              w - Inches(0.12), Inches(0.32),
              "So với tháng trước", font_name=font_b, size=9,
              color_hex=text_light, align=PP_ALIGN.CENTER)


def _text_bar(value: float, max_value: float, bars: int = 15) -> str:
    """'████████░░░░░░░' — Unicode bar tỷ lệ theo value/max_value."""
    if max_value <= 0:
        return "░" * bars
    filled = min(bars, round((value / max_value) * bars))
    return "█" * filled + "░" * (bars - filled)


def _mini_kpi(slide, x, y, w, h, label, value_str, delta_str, delta_hex, brand):
    """Mini KPI card dùng trong 2×2 grid của Zone 1."""
    font_h     = get_font(brand, "heading")
    font_b     = get_font(brand, "body")
    text_dark  = get_color(brand, "text_dark")
    text_light = get_color(brand, "text_light")
    bg_hex     = lighten_hex(get_color(brand, "primary"), 0.94)
    _add_rect(slide, x, y, w, h, bg_hex)
    stripe_col = get_color(brand, "success") if "▲" in delta_str else get_color(brand, "danger")
    _add_rect(slide, x, y, w, Pt(4), stripe_col)
    _add_text(slide, x + Inches(0.06), y + Inches(0.08),
              w - Inches(0.12), Inches(0.26),
              label, font_name=font_b, size=9,
              color_hex=text_light, align=PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.06), y + Inches(0.32),
              w - Inches(0.12), Inches(0.38),
              value_str, font_name=font_h, size=18, bold=True,
              color_hex=text_dark, align=PP_ALIGN.CENTER)
    _add_text(slide, x + Inches(0.06), y + h - Inches(0.30),
              w - Inches(0.12), Inches(0.26),
              delta_str, font_name=font_b, size=10, bold=True,
              color_hex=delta_hex, align=PP_ALIGN.CENTER)


def _mini_kpi_zone1(slide, zx, zy, zw, zh, icon, zone_title, kpis, brand):
    """Zone 1: nền + zone header + 2×2 mini KPI grid."""
    primary = get_color(brand, "primary")
    _add_rect(slide, zx, zy, zw, zh, lighten_hex(primary, 0.94))
    _zone_hdr(slide, zx, zy, zw, icon, zone_title, primary, brand)
    pad    = Inches(0.06)
    mini_w = (zw - pad * 3) / 2
    mini_h = (zh - Inches(0.4) - pad * 3) / 2
    for i, (lbl, val_str, dlt_str, dlt_hex) in enumerate(kpis[:4]):
        col = i % 2
        row = i // 2
        xi  = zx + pad + col * (mini_w + pad)
        yi  = zy + Inches(0.4) + pad + row * (mini_h + pad)
        _mini_kpi(slide, xi, yi, mini_w, mini_h, lbl, val_str, dlt_str, dlt_hex, brand)


def _data_table(slide, x, y, w, h, headers, rows, brand, *,
                header_color=None, col_widths=None, row_h_in=0.45,
                cell_text_color_fn=None, col_aligns=None):
    """Bảng dữ liệu: header màu đặc + zebra rows."""
    EMU_PER_IN = 914400
    font_h     = get_font(brand, "heading")
    font_b     = get_font(brand, "body")
    text_dark  = get_color(brand, "text_dark")
    hdr_color  = header_color or get_color(brand, "primary")
    n          = len(headers)
    if col_widths is None:
        cw_in  = (w / EMU_PER_IN) / n
        cw_list = [Inches(cw_in)] * n
    else:
        cw_list = [Inches(c) for c in col_widths]
    c_aligns   = col_aligns or ([PP_ALIGN.CENTER] + [PP_ALIGN.LEFT] * (n - 1))
    hdr_h      = Inches(0.38)
    row_h      = Inches(row_h_in)

    cx = x
    for hdr, cw, algn in zip(headers, cw_list, c_aligns):
        _add_rect(slide, cx, y, cw, hdr_h, hdr_color)
        _add_text(slide, cx + Inches(0.05), y + Inches(0.04),
                  cw - Inches(0.08), hdr_h - Inches(0.06),
                  hdr, font_name=font_h, size=9, bold=True,
                  color_hex="#FFFFFF", align=PP_ALIGN.CENTER)
        cx += cw

    max_rows = max(1, int((h - hdr_h) / row_h))
    for r, row_data in enumerate(rows[:max_rows]):
        ry     = y + hdr_h + r * row_h
        row_bg = "#F4F6FA" if r % 2 == 0 else "#FFFFFF"
        _add_rect(slide, x, ry, w, row_h, row_bg)
        cx = x
        for c, (cell, cw, algn) in enumerate(zip(row_data, cw_list, c_aligns)):
            cell_col = text_dark
            if cell_text_color_fn:
                cell_col = cell_text_color_fn(r, c, str(cell)) or text_dark
            _add_text(slide, cx + Inches(0.05), ry + Inches(0.03),
                      cw - Inches(0.08), row_h - Inches(0.05),
                      str(cell), font_name=font_b, size=10,
                      color_hex=cell_col, align=algn, wrap=True)
            cx += cw


# ---------------------------------------------------------------------------
# Đọc mock data
# ---------------------------------------------------------------------------

def load_mock_data() -> dict:
    files = {
        "gsc":        "mock_gsc.json",
        "ga4":        "mock_ga4.json",
        "ahrefs":     "mock_ahrefs.json",
        "comparison": "mock_comparison.json",
        "actions":    "mock_actions.json",
    }
    data = {}
    for key, fname in files.items():
        path = MOCK_DIR / fname
        with open(path, encoding="utf-8") as f:
            data[key] = json.load(f)
    return data


# ===========================================================================
# NHÓM A — 8 slides có nội dung cố định
# ===========================================================================

def create_slide_1_cover(prs, brand, client_name, period_str, domain, today_str):
    """Slide 1: Trang bìa."""
    slide   = _blank_slide(prs)
    primary = get_color(brand, "primary")
    accent  = get_color(brand, "accent")
    font_h  = get_font(brand, "heading")
    font_b  = get_font(brand, "body")

    _add_rect(slide, 0, 0, W, H, primary)
    _add_rect(slide, 0, H - Inches(1.0), W, Inches(1.0), accent)
    _add_rect(slide, W - Inches(4.0), 0, Inches(4.0), H - Inches(1.0), _lighten(primary))

    logo_path = get_logo_path(brand, SKILL_DIR)
    if logo_path.exists():
        slide.shapes.add_picture(
            str(logo_path),
            left=int(MARGIN), top=int(Inches(0.35)),
            width=int(Inches(2.5)), height=int(Inches(0.75))
        )

    _add_text(slide, MARGIN, Inches(1.8), W - Inches(5.0), Inches(0.55),
              "BÁO CÁO SEO HÀNG THÁNG",
              font_name=font_h, size=15, color_hex="#CCCCCC")

    _add_text(slide, MARGIN, Inches(2.35), W - Inches(5.0), Inches(1.5),
              client_name.upper(),
              font_name=font_h, size=44, bold=True, color_hex="#FFFFFF")

    _add_rect(slide, MARGIN, Inches(3.85), Inches(4.5), Pt(4), accent)

    _add_text(slide, MARGIN, Inches(4.0), W - Inches(5.0), Inches(0.7),
              period_str,
              font_name=font_h, size=24, bold=True, color_hex=accent)

    _add_text(slide, MARGIN, Inches(4.7), W - Inches(5.0), Inches(0.5),
              domain, font_name=font_b, size=14, color_hex="#AAAAAA")

    _add_text(slide, MARGIN, H - Inches(0.9), W - 2*MARGIN, Inches(0.5),
              f"Báo cáo được chuẩn bị bởi SEONGON Agency  •  {today_str}",
              font_name=font_b, size=11, color_hex="#FFFFFF",
              align=PP_ALIGN.CENTER)


def create_slide_2_toc(prs, brand):
    """Slide 2: Mục lục."""
    slide      = _blank_slide(prs)
    primary    = get_color(brand, "primary")
    accent     = get_color(brand, "accent")
    text_dark  = get_color(brand, "text_dark")
    text_light = get_color(brand, "text_light")
    font_h     = get_font(brand, "heading")
    font_b     = get_font(brand, "body")

    _add_rect(slide, 0, 0, W, H, get_color(brand, "background"))
    _add_rect(slide, 0, 0, Inches(0.2), H, primary)
    _add_text(slide, MARGIN, Inches(0.38), Inches(8), Inches(0.85),
              "Mục lục", font_name=font_h, size=34, bold=True, color_hex=primary)
    _add_rect(slide, MARGIN, Inches(1.22), Inches(4.5), Pt(3), accent)

    sections = [
        ("I",   "Tóm tắt điều hành",       "Slide 3",        "Điểm sáng & điểm cần chú ý trong tháng"),
        ("II",  "Hiệu suất tổng quan",      "Slides 4 – 6",   "KPI chính, xu hướng lưu lượng & chuyển đổi"),
        ("III", "Phân tích chi tiết",       "Slides 7 – 10",  "GA4 • GSC • Backlinks — phân tích 4 góc nhìn"),
        ("IV",  "Kế hoạch hành động",       "Slides 11 – 13", "Việc đã hoàn thành & kế hoạch tháng tới"),
        ("V",   "Khuyến nghị chiến lược",   "Slide 14",       "Định hướng dài hạn Q3–Q4/2026 và 2027"),
    ]

    start_y = Inches(1.45)
    row_h   = Inches(1.0)
    for i, (num, title, slide_ref, desc) in enumerate(sections):
        y = start_y + i * row_h
        if i % 2 == 0:
            _add_rect(slide, Inches(0.28), y, W - Inches(0.28),
                      row_h - Inches(0.04), "#F4F6FA")
        _add_rect(slide, Inches(0.36), y + Inches(0.18),
                  Inches(0.52), Inches(0.52), primary)
        _add_text(slide, Inches(0.36), y + Inches(0.15),
                  Inches(0.52), Inches(0.52), num,
                  font_name=font_h, size=10, bold=True,
                  color_hex="#FFFFFF", align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(1.05), y + Inches(0.1),
                  Inches(6.5), Inches(0.42),
                  title, font_name=font_h, size=15, bold=True, color_hex=text_dark)
        _add_text(slide, Inches(1.05), y + Inches(0.52),
                  Inches(8.5), Inches(0.38),
                  desc, font_name=font_b, size=11, color_hex=text_light)
        _add_text(slide, W - Inches(2.3), y + Inches(0.23),
                  Inches(1.95), Inches(0.4), slide_ref,
                  font_name=font_b, size=11, bold=True,
                  color_hex=primary, align=PP_ALIGN.RIGHT)

    _footer(slide, brand)


def create_slide_3_exec_summary(prs, brand, comparison_data):
    """Slide 3: Tóm tắt điều hành."""
    slide      = _blank_slide(prs)
    primary    = get_color(brand, "primary")
    accent     = get_color(brand, "accent")
    bg         = get_color(brand, "background")
    text_dark  = get_color(brand, "text_dark")
    text_light = get_color(brand, "text_light")
    success    = get_color(brand, "success")
    danger     = get_color(brand, "danger")
    warning    = get_color(brand, "warning")
    font_h     = get_font(brand, "heading")
    font_b     = get_font(brand, "body")

    _add_rect(slide, 0, 0, W, H, bg)
    _add_rect(slide, 0, 0, Inches(0.2), H, primary)
    _add_text(slide, MARGIN, Inches(0.35), Inches(9), Inches(0.75),
              "Tóm tắt điều hành",
              font_name=font_h, size=30, bold=True, color_hex=primary)
    _add_rect(slide, MARGIN, Inches(1.1), Inches(5), Pt(3), accent)

    overall   = comparison_data.get("overall_assessment", {})
    score     = overall.get("score", 0)
    max_score = overall.get("max_score", 10)
    label_oa  = overall.get("label", "")
    _add_rect(slide, W - Inches(3.8), Inches(0.28), Inches(3.3), Inches(0.72), primary)
    _add_text(slide, W - Inches(3.8), Inches(0.28), Inches(3.3), Inches(0.72),
              f"Điểm tháng: {score}/{max_score}  •  {label_oa}",
              font_name=font_b, size=10, bold=True,
              color_hex="#FFFFFF", align=PP_ALIGN.CENTER)

    metrics      = comparison_data.get("metrics", [])
    good_items   = [m for m in metrics if m.get("status") == "good"][:4]
    bad_items    = [m for m in metrics if m.get("status") == "bad"][:3]
    neutral_items= [m for m in metrics if m.get("status") == "neutral"][:1]

    col_w  = (W - 2*MARGIN - Inches(0.35)) / 2
    col1_x = MARGIN
    col2_x = MARGIN + col_w + Inches(0.35)
    top_y  = Inches(1.25)
    col_h  = Inches(4.3)

    _add_rect(slide, col1_x, top_y, col_w, col_h, "#EFF8F1")
    _add_rect(slide, col1_x, top_y, col_w, Inches(0.44), success)
    _add_text(slide, col1_x + Inches(0.12), top_y + Inches(0.05),
              col_w - Inches(0.15), Inches(0.36),
              "✅  Điểm sáng tháng này",
              font_name=font_h, size=12, bold=True, color_hex="#FFFFFF")
    by = top_y + Inches(0.55)
    for m in good_items:
        dp   = m.get("delta_pct", 0)
        sign = "+" if dp >= 0 else ""
        pct  = f"{sign}{dp*100:.1f}%".replace(".", ",")
        name = m.get("name", "").split("(")[0].strip()[:28]
        _add_text(slide, col1_x + Inches(0.15), by, col_w - Inches(0.2), Inches(0.62),
                  f"▲  {name}: {pct} so với tháng trước",
                  font_name=font_b, size=11, color_hex=text_dark, wrap=True)
        by += Inches(0.82)

    _add_rect(slide, col2_x, top_y, col_w, col_h, "#FFF5F0")
    _add_rect(slide, col2_x, top_y, col_w, Inches(0.44), danger)
    _add_text(slide, col2_x + Inches(0.12), top_y + Inches(0.05),
              col_w - Inches(0.15), Inches(0.36),
              "⚠️  Điểm cần chú ý",
              font_name=font_h, size=12, bold=True, color_hex="#FFFFFF")
    by2 = top_y + Inches(0.55)
    for m in bad_items:
        dp   = m.get("delta_pct", 0)
        sign = "+" if dp >= 0 else ""
        pct  = f"{sign}{dp*100:.1f}%".replace(".", ",")
        name = m.get("name", "").split("(")[0].strip()[:28]
        _add_text(slide, col2_x + Inches(0.15), by2, col_w - Inches(0.2), Inches(0.62),
                  f"▼  {name}: {pct} — cần theo dõi",
                  font_name=font_b, size=11, color_hex=text_dark, wrap=True)
        by2 += Inches(0.82)
    for m in neutral_items:
        name = m.get("name", "").split("(")[0].strip()[:28]
        _add_text(slide, col2_x + Inches(0.15), by2, col_w - Inches(0.2), Inches(0.62),
                  f"→  {name}: Ổn định, cần theo dõi",
                  font_name=font_b, size=11, color_hex=text_dark, wrap=True)

    summary_text = overall.get("summary", "")
    _add_rect(slide, MARGIN, H - Inches(1.35), W - 2*MARGIN, Inches(0.95), "#ECF1F8")
    _add_text(slide, MARGIN + Inches(0.14), H - Inches(1.3),
              W - 2*MARGIN - Inches(0.2), Inches(0.88),
              f"📋  {summary_text}",
              font_name=font_b, size=10, color_hex=text_dark, wrap=True)

    _footer(slide, brand)


def create_slide_divider(prs, brand, roman_num: str, title: str, subtitle: str = ""):
    """Section divider — dùng lại cho slides 4, 7, 11."""
    slide   = _blank_slide(prs)
    primary = get_color(brand, "primary")
    accent  = get_color(brand, "accent")
    font_h  = get_font(brand, "heading")
    font_b  = get_font(brand, "body")

    _add_rect(slide, 0, 0, W, H, primary)
    _add_rect(slide, 0, H - Inches(0.42), W, Inches(0.42), accent)
    _add_rect(slide, W - Inches(4.0), 0, Inches(4.0), H - Inches(0.42), _lighten(primary))

    _add_text(slide, W - Inches(4.2), Inches(0.5), Inches(4.0), Inches(5.5),
              roman_num, font_name=font_h, size=180, bold=True,
              color_hex="#FFFFFF", align=PP_ALIGN.CENTER, wrap=False)

    _add_text(slide, MARGIN, Inches(2.4), W - Inches(5.0), Inches(1.4),
              title.upper(), font_name=font_h, size=36, bold=True, color_hex="#FFFFFF")

    _add_rect(slide, MARGIN, Inches(3.85), Inches(5.0), Pt(4), accent)

    if subtitle:
        _add_text(slide, MARGIN, Inches(4.05), W - Inches(5.0), Inches(0.6),
                  subtitle, font_name=font_b, size=15, color_hex="#BBBBBB")

    _add_text(slide, MARGIN, H - Inches(0.38), Inches(6), Inches(0.35),
              "seongon.com  •  Báo cáo SEO",
              font_name=font_b, size=10, color_hex="#FFFFFF")


def create_slide_14_longterm(prs, brand, actions_data):
    """Slide 14: Khuyến nghị chiến lược dài hạn."""
    slide      = _blank_slide(prs)
    primary    = get_color(brand, "primary")
    accent     = get_color(brand, "accent")
    bg         = get_color(brand, "background")
    text_dark  = get_color(brand, "text_dark")
    text_light = get_color(brand, "text_light")
    success    = get_color(brand, "success")
    warning    = get_color(brand, "warning")
    font_h     = get_font(brand, "heading")
    font_b     = get_font(brand, "body")

    _add_rect(slide, 0, 0, W, H, bg)
    _add_rect(slide, 0, 0, Inches(0.2), H, primary)
    _add_text(slide, MARGIN, Inches(0.35), Inches(9), Inches(0.75),
              "Khuyến nghị chiến lược dài hạn",
              font_name=font_h, size=28, bold=True, color_hex=primary)
    _add_rect(slide, MARGIN, Inches(1.1), Inches(5.5), Pt(3), accent)

    long_term      = actions_data.get("long_term_recommendations", [])[:5]
    priority_color = {"Cao": success, "Trung bình": warning, "Thấp": text_light}
    icons          = ["🎯", "🏆", "📍", "📹", "📊"]
    row_h          = Inches(1.15)
    start_y        = Inches(1.28)

    for i, rec in enumerate(long_term):
        y       = start_y + i * row_h
        pri     = rec.get("priority", "Thấp")
        p_color = priority_color.get(pri, text_light)
        icon    = icons[i % len(icons)]
        if i % 2 == 0:
            _add_rect(slide, MARGIN, y, W - 2*MARGIN, row_h - Inches(0.04), "#F4F6FA")
        badge = {"Cao": "C", "Trung bình": "TB", "Thấp": "T"}.get(pri, pri[0])
        _add_rect(slide, MARGIN, y + Inches(0.15), Inches(0.55), Inches(0.55), p_color)
        _add_text(slide, MARGIN, y + Inches(0.13), Inches(0.55), Inches(0.55),
                  badge, font_name=font_h, size=10, bold=True,
                  color_hex="#FFFFFF", align=PP_ALIGN.CENTER)
        cx = MARGIN + Inches(0.72)
        cw = W - 2*MARGIN - Inches(0.72) - Inches(2.1)
        _add_text(slide, cx, y + Inches(0.09), cw, Inches(0.4),
                  f"{icon}  {rec.get('recommendation', '')}",
                  font_name=font_h, size=12, bold=True, color_hex=text_dark)
        detail = f"⏱ {rec.get('timeline', '')}  •  Dự kiến: {rec.get('expected_impact', '')[:65]}"
        _add_text(slide, cx, y + Inches(0.52), cw, Inches(0.5),
                  detail, font_name=font_b, size=10, color_hex=text_light)
        _add_text(slide, W - MARGIN - Inches(1.9), y + Inches(0.22),
                  Inches(1.8), Inches(0.38),
                  f"Ưu tiên: {pri}", font_name=font_b, size=10,
                  bold=True, color_hex=p_color, align=PP_ALIGN.RIGHT)

    _footer(slide, brand)


def create_slide_15_thankyou(prs, brand, domain):
    """Slide 15: Cảm ơn / Liên hệ."""
    slide   = _blank_slide(prs)
    primary = get_color(brand, "primary")
    accent  = get_color(brand, "accent")
    font_h  = get_font(brand, "heading")
    font_b  = get_font(brand, "body")

    _add_rect(slide, 0, 0, W, H, primary)
    _add_rect(slide, 0, H - Inches(0.55), W, Inches(0.55), accent)
    _add_rect(slide, W - Inches(3.0), -Inches(0.5), Inches(3.0), Inches(3.0), _lighten(primary, 0.08))
    _add_rect(slide, W - Inches(1.8), Inches(0.8), Inches(2.2), Inches(2.2), _lighten(primary, 0.15))

    _add_text(slide, MARGIN, Inches(1.85), W - 2*MARGIN, Inches(1.2),
              "Cảm ơn quý khách!",
              font_name=font_h, size=52, bold=True,
              color_hex="#FFFFFF", align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(3.1), W - 2*MARGIN, Inches(0.6),
              "Chúng tôi cam kết đồng hành cùng sự phát triển của doanh nghiệp bạn.",
              font_name=font_b, size=14, color_hex="#CCCCCC",
              align=PP_ALIGN.CENTER)
    _add_rect(slide, int((W - Inches(4.0)) / 2), Inches(3.85),
              Inches(4.0), Pt(2), accent)

    contact_items = [
        ("🌐", "seongon.com"),
        ("📧", "contact@seongon.com"),
        ("📞", "(028) 7109 xxxx"),
    ]
    item_w  = Inches(3.8)
    start_x = (W - item_w * len(contact_items)) / 2
    for j, (icon, info) in enumerate(contact_items):
        cx = start_x + j * item_w
        _add_text(slide, cx, Inches(4.1), item_w, Inches(0.5),
                  f"{icon}  {info}",
                  font_name=font_b, size=14, color_hex="#FFFFFF",
                  align=PP_ALIGN.CENTER)

    logo_path = get_logo_path(brand, SKILL_DIR)
    if logo_path.exists():
        lw = Inches(2.5)
        lh = Inches(0.75)
        slide.shapes.add_picture(str(logo_path),
            left=int((W - lw) / 2), top=int(H - Inches(1.6)),
            width=int(lw), height=int(lh))

    _add_text(slide, MARGIN, H - Inches(0.52), W - 2*MARGIN, Inches(0.45),
              f"© 2026 SEONGON Agency  •  Tài liệu bảo mật — chỉ dành cho {brand['client_name']}",
              font_name=font_b, size=9, color_hex="#FFFFFF",
              align=PP_ALIGN.CENTER)


# ===========================================================================
# NHÓM B — Sub-step B1: Slides 5 và 6
# ===========================================================================

def create_slide_5_kpi_dashboard(prs, brand, ga4_data, gsc_data):
    """Slide 5: KPI Summary Dashboard — 2×3 grid."""
    slide = _blank_slide(prs)
    _slide_base(slide, brand, "KPI Summary Dashboard")

    ga4c = ga4_data["summary"]["current"]
    ga4p = ga4_data["summary"]["previous"]
    gscc = gsc_data["summary"]["current"]
    gscp = gsc_data["summary"]["previous"]
    text_dark  = get_color(brand, "text_dark")
    font_b     = get_font(brand, "body")

    # 6 KPI: label, current, previous, higher_is_better, format_fn
    def fmt_pos(v):
        return f"{float(v):.1f}".replace(".", ",")

    kpis = [
        ("Phiên truy cập (Sessions)", ga4c["sessions"],       ga4p["sessions"],       True,  format_vn_number),
        ("Người dùng (Users)",         ga4c["users"],          ga4p["users"],          True,  format_vn_number),
        ("Chuyển đổi (Conversions)",   ga4c["conversions"],    ga4p["conversions"],    True,  format_vn_number),
        ("Lượt nhấp (Clicks GSC)",     gscc["clicks"],         gscp["clicks"],         True,  format_vn_number),
        ("Hiển thị (Impressions GSC)", gscc["impressions"],    gscp["impressions"],    True,  format_vn_number),
        ("Vị trí TB (Avg Position)",   gscc["avg_position"],   gscp["avg_position"],   False, fmt_pos),
    ]

    card_w = Inches(3.978)
    card_h = Inches(2.35)
    gap_x  = Inches(0.2)
    gap_y  = Inches(0.18)
    row1_y = Inches(1.2)
    row2_y = row1_y + card_h + gap_y
    xs     = [MARGIN, MARGIN + card_w + gap_x, MARGIN + 2 * (card_w + gap_x)]

    good_count = 0
    for i, (label, cur, prev, hib, fmt) in enumerate(kpis):
        row = i // 3
        col = i % 3
        x   = xs[col]
        y   = row1_y if row == 0 else row2_y

        val_str             = fmt(cur)
        delta_str, d_color  = _calc_delta(cur, prev, brand, hib)

        # Màu top-stripe theo tốt/xấu
        try:
            d = (float(cur) - float(prev)) / float(prev)
            is_good = (d >= 0) if hib else (d <= 0)
            if is_good:
                good_count += 1
        except Exception:
            is_good = True
        top_color = get_color(brand, "success") if is_good else get_color(brand, "danger")

        _kpi_card(slide, x, y, card_w, card_h, label,
                  val_str, delta_str, d_color, brand, top_bar_hex=top_color)

    # Summary bar
    bad_count = len(kpis) - good_count
    _add_rect(slide, MARGIN, Inches(6.12), W - 2*MARGIN, Inches(0.65), "#EEF2F7")
    _add_text(slide, MARGIN + Inches(0.2), Inches(6.17),
              W - 2*MARGIN - Inches(0.3), Inches(0.55),
              f"📊  Tổng quan: {good_count}/6 KPI tăng tích cực  •  {bad_count}/6 cần theo dõi",
              font_name=font_b, size=13, bold=True,
              color_hex=text_dark, align=PP_ALIGN.CENTER)

    _footer(slide, brand)


def create_slide_6_traffic_trend(prs, brand, ga4_data, gsc_data, comparison_data):
    """Slide 6: Xu hướng lưu lượng & chuyển đổi — layout 4-zone 2×2."""
    slide      = _blank_slide(prs)
    _slide_base(slide, brand, "Xu hướng Lưu lượng & Chuyển đổi")

    primary    = get_color(brand, "primary")
    success    = get_color(brand, "success")
    warning    = get_color(brand, "warning")
    text_dark  = get_color(brand, "text_dark")
    text_light = get_color(brand, "text_light")
    font_h     = get_font(brand, "heading")
    font_b     = get_font(brand, "body")

    ga4c = ga4_data["summary"]["current"]
    ga4p = ga4_data["summary"]["previous"]

    # --- ZONE 1 (top-left): Sessions KPI lớn ---
    z1_bg = lighten_hex(primary, 0.92)
    _add_rect(slide, MARGIN, _ZTY, _ZLW, _ZTH, z1_bg)
    _zone_hdr(slide, MARGIN, _ZTY, _ZLW, "📈", "Tăng trưởng Lưu lượng", primary, brand)

    sessions_cur  = ga4c["sessions"]
    sessions_prev = ga4p["sessions"]
    delta_str, d_color = _calc_delta(sessions_cur, sessions_prev, brand, True)

    _add_text(slide, MARGIN + Inches(0.2), _ZTY + Inches(0.55),
              _ZLW - Inches(0.3), Inches(1.0),
              format_vn_number(sessions_cur),
              font_name=font_h, size=50, bold=True,
              color_hex=primary, align=PP_ALIGN.CENTER)

    _add_text(slide, MARGIN + Inches(0.2), _ZTY + Inches(1.6),
              _ZLW - Inches(0.3), Inches(0.38),
              "phiên truy cập  —  Tháng 4/2026",
              font_name=font_b, size=11, color_hex=text_light,
              align=PP_ALIGN.CENTER)

    _add_text(slide, MARGIN + Inches(0.2), _ZTY + Inches(2.05),
              _ZLW - Inches(0.3), Inches(0.42),
              f"{delta_str}  so với tháng 3/2026",
              font_name=font_h, size=16, bold=True,
              color_hex=d_color, align=PP_ALIGN.CENTER)

    # --- ZONE 2 (top-right): Bar chart lượt nhấp GSC theo tuần ---
    _add_rect(slide, _ZRX, _ZTY, _ZLW, _ZTH, "#F8F9FA")
    _zone_hdr(slide, _ZRX, _ZTY, _ZLW, "📊", "Clicks GSC theo tuần (W14 – W17)", primary, brand)

    gsc_weekly = gsc_data.get("weekly_trend", [])
    if gsc_weekly:
        max_clicks = max(w["clicks"] for w in gsc_weekly)
        cy         = _ZTY + Inches(0.52)
        for w_item in gsc_weekly:
            bar     = _text_bar(w_item["clicks"], max_clicks, 14)
            row_txt = f"{w_item['week']}  {bar}  {format_vn_number(w_item['clicks'])}"
            _add_text(slide, _ZRX + Inches(0.15), cy,
                      _ZLW - Inches(0.25), Inches(0.5),
                      row_txt, font_name="Courier New", size=12,
                      color_hex=primary)
            cy += Inches(0.52)

    # --- ZONE 3 (bottom-left): Điểm tốt ---
    _add_rect(slide, MARGIN, _ZBY, _ZLW, _ZBH, "#EFF8F1")
    _zone_hdr(slide, MARGIN, _ZBY, _ZLW, "✅", "Điểm tốt", success, brand)

    sessions_delta_pct = (sessions_cur - sessions_prev) / sessions_prev * 100
    conv_cur  = ga4c["conversions"]
    conv_prev = ga4p["conversions"]

    good_bullets = [
        f"Sessions tăng {sessions_delta_pct:.1f}% vs T3".replace(".", ",")
        + " — vượt mục tiêu +10% đề ra",
        f"Conversions tăng mạnh: {format_vn_number(conv_cur)} "
        f"(+{(conv_cur - conv_prev) / conv_prev * 100:.1f}%".replace(".", ",") + " so T3)",
        "Non-branded keywords tăng 16,3% — tăng trưởng chủ động",
        "Organic chiếm 86,9% tổng sessions — tỷ trọng rất cao",
    ]
    by = _ZBY + Inches(0.5)
    for bullet in good_bullets:
        _add_text(slide, MARGIN + Inches(0.15), by,
                  _ZLW - Inches(0.22), Inches(0.48),
                  f"•  {bullet}", font_name=font_b, size=11,
                  color_hex=text_dark, wrap=True)
        by += Inches(0.5)

    # --- ZONE 4 (bottom-right): Điểm chưa tốt ---
    _add_rect(slide, _ZRX, _ZBY, _ZLW, _ZBH, "#FFFBF0")
    _zone_hdr(slide, _ZRX, _ZBY, _ZLW, "⚠️", "Điểm chưa tốt", warning, brand)

    bounce_cur  = ga4c.get("bounce_rate", 0) * 100
    bounce_prev = ga4p.get("bounce_rate", 0) * 100
    dur_cur     = ga4c.get("avg_session_duration_s", 0)
    dur_prev    = ga4p.get("avg_session_duration_s", 0)

    bad_bullets = [
        f"Bounce rate tăng: {bounce_prev:.1f}% → {bounce_cur:.1f}%".replace(".", ",")
        + f" (+{bounce_cur - bounce_prev:.1f}pp)".replace(".", ","),
        f"Thời gian phiên TB giảm: {dur_prev}s → {dur_cur}s (−14,6%)",
        "CTR giữ nguyên 6,21% dù avg position cải thiện",
        "Cần tối ưu Core Web Vitals để giảm bounce rate",
    ]
    by2 = _ZBY + Inches(0.5)
    for bullet in bad_bullets:
        _add_text(slide, _ZRX + Inches(0.15), by2,
                  _ZLW - Inches(0.22), Inches(0.48),
                  f"•  {bullet}", font_name=font_b, size=11,
                  color_hex=text_dark, wrap=True)
        by2 += Inches(0.5)

    # --- Narrative ---
    metrics_by_id = {m["id"]: m for m in comparison_data.get("metrics", [])}
    narr_text = metrics_by_id.get("sessions", {}).get(
        "interpretation",
        "Lưu lượng tháng 4 tăng mạnh 15,1% nhờ link building T3 phát huy. "
        "Ưu tiên T5: cải thiện bounce rate và thời gian phiên."
    )[:130]
    _narr(slide, brand, narr_text)
    _footer(slide, brand)


# ===========================================================================
# NHÓM B — Sub-step B2: Slides 8, 9, 10, 12, 13
# ===========================================================================

def create_slide_8_ga4(prs, brand, ga4_data, comparison_data):
    """Slide 8: Phân tích GA4 — layout 4-zone 2×2."""
    slide = _blank_slide(prs)
    _slide_base(slide, brand, "Phân tích GA4 — Hành vi người dùng")

    primary    = get_color(brand, "primary")
    success    = get_color(brand, "success")
    warning    = get_color(brand, "warning")
    text_dark  = get_color(brand, "text_dark")
    font_b     = get_font(brand, "body")

    ga4c = ga4_data["summary"]["current"]
    ga4p = ga4_data["summary"]["previous"]

    # --- ZONE 1 (top-left): 2×2 mini KPI ---
    def _fmt_br(v):
        return f"{float(v)*100:.1f}%".replace(".", ",")

    kpis_z1 = [
        ("Sessions",    format_vn_number(ga4c["sessions"]),
         *_calc_delta(ga4c["sessions"],    ga4p["sessions"],    brand, True)),
        ("Users",       format_vn_number(ga4c["users"]),
         *_calc_delta(ga4c["users"],       ga4p["users"],       brand, True)),
        ("Conversions", format_vn_number(ga4c["conversions"]),
         *_calc_delta(ga4c["conversions"], ga4p["conversions"], brand, True)),
        ("Bounce Rate", _fmt_br(ga4c["bounce_rate"]),
         *_calc_delta(ga4c["bounce_rate"], ga4p["bounce_rate"], brand, False)),
    ]
    _mini_kpi_zone1(slide, MARGIN, _ZTY, _ZLW, _ZTH,
                    "📊", "Chỉ số tổng quan GA4", kpis_z1, brand)

    # --- ZONE 2 (top-right): Top 5 trang đích ---
    _add_rect(slide, _ZRX, _ZTY, _ZLW, _ZTH, "#F8F9FA")
    _zone_hdr(slide, _ZRX, _ZTY, _ZLW, "📄", "Top 5 Trang đích (Sessions)", primary, brand)
    pages    = ga4_data["top_landing_pages"][:5]
    tbl_rows = []
    for pg in pages:
        try:
            pct  = (pg["sessions"] - pg["prev_sessions"]) / pg["prev_sessions"] * 100
            sign = "+" if pct >= 0 else ""
            d_s  = f"{sign}{pct:.1f}%".replace(".", ",")
        except Exception:
            d_s = "N/A"
        tbl_rows.append([pg["page"][:24], format_vn_number(pg["sessions"]), d_s])
    _data_table(slide, _ZRX, _ZTY + Inches(0.4),
                _ZLW, _ZTH - Inches(0.4),
                ["Trang đích", "Sessions", "Δ%"], tbl_rows, brand,
                header_color=primary, col_widths=[3.567, 1.35, 1.15], row_h_in=0.43,
                col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])

    # --- ZONE 3 (bottom-left ✅): Điểm tốt ---
    _add_rect(slide, MARGIN, _ZBY, _ZLW, _ZBH, "#EFF8F1")
    _zone_hdr(slide, MARGIN, _ZBY, _ZLW, "✅", "Điểm tốt — GA4", success, brand)
    bounce_c = ga4c["bounce_rate"] * 100
    bounce_p = ga4p["bounce_rate"] * 100
    dur_c    = ga4c["avg_session_duration_s"]
    dur_p    = ga4p["avg_session_duration_s"]
    good_bullets = [
        f"Sessions: {format_vn_number(ga4p['sessions'])} → {format_vn_number(ga4c['sessions'])} (+15,1%) — vượt mục tiêu +10%",
        f"Conversions: {format_vn_number(ga4p['conversions'])} → {format_vn_number(ga4c['conversions'])} (+22,4%) — tốc độ tăng cao hơn sessions",
        "Non-branded sessions tăng 16,3% — tăng trưởng chủ động, không phụ thuộc thương hiệu",
        "Organic traffic chiếm 86,9% tổng sessions — tỷ trọng rất cao, bền vững",
    ]
    by = _ZBY + Inches(0.5)
    for b in good_bullets:
        _add_text(slide, MARGIN + Inches(0.15), by, _ZLW - Inches(0.22), Inches(0.5),
                  f"•  {b}", font_name=font_b, size=10, color_hex=text_dark, wrap=True)
        by += Inches(0.5)

    # --- ZONE 4 (bottom-right ⚠️): Điểm chưa tốt ---
    _add_rect(slide, _ZRX, _ZBY, _ZLW, _ZBH, "#FFFBF0")
    _zone_hdr(slide, _ZRX, _ZBY, _ZLW, "⚠️", "Điểm chưa tốt — GA4", warning, brand)
    bad_bullets = [
        f"Bounce rate tăng: {bounce_p:.1f}% → {bounce_c:.1f}% (+{bounce_c - bounce_p:.1f}pp) — cần xử lý ngay".replace(".", ","),
        f"Thời gian phiên TB giảm: {dur_p}s → {dur_c}s (−14,6%) — nội dung chưa đủ engaging",
        "Trang /kien-thuc/seo-la-gi: 1.950 sessions nhưng chỉ 8 conversions (0,4%)",
        "Ưu tiên T5: tối ưu Core Web Vitals + thêm CTA vào bài blog traffic cao",
    ]
    by2 = _ZBY + Inches(0.5)
    for b in bad_bullets:
        _add_text(slide, _ZRX + Inches(0.15), by2, _ZLW - Inches(0.22), Inches(0.5),
                  f"•  {b}", font_name=font_b, size=10, color_hex=text_dark, wrap=True)
        by2 += Inches(0.5)

    _narr(slide, brand,
          "GA4 T4: Sessions +15,1%, Conversions +22,4% — tích cực. "
          "Cần ưu tiên giảm bounce rate (+3,8pp) và cải thiện thời gian phiên (−14,6%) trong T5.")
    _footer(slide, brand)


def create_slide_9_gsc(prs, brand, gsc_data, comparison_data):
    """Slide 9: Phân tích GSC — layout 4-zone 2×2."""
    slide = _blank_slide(prs)
    _slide_base(slide, brand, "Phân tích GSC — Hiệu suất tìm kiếm")

    primary    = get_color(brand, "primary")
    success    = get_color(brand, "success")
    warning    = get_color(brand, "warning")
    text_dark  = get_color(brand, "text_dark")
    font_b     = get_font(brand, "body")

    gscc   = gsc_data["summary"]["current"]
    gscp   = gsc_data["summary"]["previous"]
    kw_cur = gsc_data["organic_keywords"]["current"]
    kw_prv = gsc_data["organic_keywords"]["previous"]

    # --- ZONE 1 (top-left): 2×2 mini KPI ---
    def _fmt_ctr(v):
        return f"{float(v)*100:.2f}%".replace(".", ",")
    def _fmt_pos(v):
        return f"{float(v):.1f}".replace(".", ",")

    kpis_z1 = [
        ("Clicks",         format_vn_number(gscc["clicks"]),
         *_calc_delta(gscc["clicks"],       gscp["clicks"],       brand, True)),
        ("Impressions",    format_vn_number(gscc["impressions"]),
         *_calc_delta(gscc["impressions"],  gscp["impressions"],  brand, True)),
        ("CTR trung bình", _fmt_ctr(gscc["ctr"]),
         *_calc_delta(gscc["ctr"],          gscp["ctr"],          brand, True)),
        ("Avg Position",   _fmt_pos(gscc["avg_position"]),
         *_calc_delta(gscc["avg_position"], gscp["avg_position"], brand, False)),
    ]
    _mini_kpi_zone1(slide, MARGIN, _ZTY, _ZLW, _ZTH,
                    "🔍", "Chỉ số tổng quan GSC", kpis_z1, brand)

    # --- ZONE 2 (top-right): Từ khóa biến động ---
    _add_rect(slide, _ZRX, _ZTY, _ZLW, _ZTH, "#F8F9FA")
    _zone_hdr(slide, _ZRX, _ZTY, _ZLW, "📈", "Từ khóa biến động mạnh", primary, brand)
    kw_changes = gsc_data.get("keywords_with_big_changes", [])
    kw_rows    = []
    for kw in kw_changes[:7]:
        chg = abs(kw["position_change"])
        if kw["direction"] == "improved":
            arrow   = "↑"
            chg_str = f"▲ {chg:.1f} bậc".replace(".", ",")
        else:
            arrow   = "↓"
            chg_str = f"▼ +{chg:.1f} bậc".replace(".", ",")
        kw_rows.append([kw["query"][:22], arrow, chg_str])

    def _kw_color(r, c, v):
        if c == 1:
            return get_color(brand, "success") if v == "↑" else get_color(brand, "danger")
        return None

    _data_table(slide, _ZRX, _ZTY + Inches(0.4),
                _ZLW, _ZTH - Inches(0.4),
                ["Từ khóa", "↑↓", "Thay đổi vị trí"], kw_rows, brand,
                header_color=primary, col_widths=[3.567, 0.65, 1.85], row_h_in=0.32,
                col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT],
                cell_text_color_fn=_kw_color)

    # --- ZONE 3 (bottom-left ✅): Điểm tốt ---
    _add_rect(slide, MARGIN, _ZBY, _ZLW, _ZBH, "#EFF8F1")
    _zone_hdr(slide, MARGIN, _ZBY, _ZLW, "✅", "Điểm tốt — GSC", success, brand)
    good_bullets = [
        f"Clicks tăng 11,3%: {format_vn_number(gscp['clicks'])} → {format_vn_number(gscc['clicks'])} — nhờ thứ hạng cải thiện",
        f"Avg Position cải thiện: {str(gscp['avg_position']).replace('.', ',')} → {str(gscc['avg_position']).replace('.', ',')} (−0,8 bậc)",
        f"Keywords top 10: {kw_prv['top_10']} → {kw_cur['top_10']} (+23 từ khóa mới vào top 10)",
        f"Keywords top 3: {kw_prv['top_3']} → {kw_cur['top_3']} (+7) — thứ hạng cao nhất tăng mạnh",
    ]
    by = _ZBY + Inches(0.5)
    for b in good_bullets:
        _add_text(slide, MARGIN + Inches(0.15), by, _ZLW - Inches(0.22), Inches(0.5),
                  f"•  {b}", font_name=font_b, size=10, color_hex=text_dark, wrap=True)
        by += Inches(0.5)

    # --- ZONE 4 (bottom-right ⚠️): Điểm chưa tốt ---
    _add_rect(slide, _ZRX, _ZBY, _ZLW, _ZBH, "#FFFBF0")
    _zone_hdr(slide, _ZRX, _ZBY, _ZLW, "⚠️", "Điểm chưa tốt — GSC", warning, brand)
    ctr_c = gscc["ctr"] * 100
    ctr_p = gscp["ctr"] * 100
    bad_bullets = [
        f"CTR gần như không đổi: {ctr_p:.2f}% → {ctr_c:.2f}% — position cải thiện nhưng CTR không tăng".replace(".", ","),
        "8 từ khóa giảm hạng: báo giá seo, seo tổng thể, kỹ thuật seo — cần cập nhật nội dung",
        "3 trang giảm clicks: /giai-phap-seo (−15,4%), /ve-chung-toi (−16,4%)",
        "Ưu tiên T5: tối ưu meta title/description cho 20 trang CTR < 4%, impression > 500",
    ]
    by2 = _ZBY + Inches(0.5)
    for b in bad_bullets:
        _add_text(slide, _ZRX + Inches(0.15), by2, _ZLW - Inches(0.22), Inches(0.5),
                  f"•  {b}", font_name=font_b, size=10, color_hex=text_dark, wrap=True)
        by2 += Inches(0.5)

    _narr(slide, brand,
          "GSC T4: Clicks +11,3%, Avg Position 9,1 → 8,3 — tích cực. "
          "CTR giữ nguyên 6,21% dù position cải thiện — cần tối ưu meta data T5.")
    _footer(slide, brand)


def create_slide_10_backlinks(prs, brand, ahrefs_data, comparison_data):
    """Slide 10: Phân tích Backlinks & Domain Rating — layout 4-zone 2×2."""
    slide = _blank_slide(prs)
    _slide_base(slide, brand, "Phân tích Backlinks & Domain Rating")

    primary    = get_color(brand, "primary")
    success    = get_color(brand, "success")
    warning    = get_color(brand, "warning")
    text_dark  = get_color(brand, "text_dark")
    text_light = get_color(brand, "text_light")
    font_h     = get_font(brand, "heading")
    font_b     = get_font(brand, "body")

    dr_cur  = ahrefs_data["domain_rating"]["current"]
    dr_prev = ahrefs_data["domain_rating"]["previous"]
    rd_cur  = ahrefs_data["referring_domains"]["current"]
    rd_prev = ahrefs_data["referring_domains"]["previous"]
    rd_new  = ahrefs_data["referring_domains"]["new_this_month"]
    rd_lost = ahrefs_data["referring_domains"]["lost_this_month"]
    bl_tot  = ahrefs_data["backlinks"]["total"]
    bl_prev = ahrefs_data["backlinks"]["total_previous"]

    # --- ZONE 1 (top-left): 2×2 mini KPI ---
    kpis_z1 = [
        ("Domain Rating (DR)", str(dr_cur),
         *_calc_delta(dr_cur, dr_prev, brand, True, pct=False)),
        ("Referring Domains",  format_vn_number(rd_cur),
         *_calc_delta(rd_cur, rd_prev, brand, True)),
        ("Tổng Backlinks",     format_vn_number(bl_tot),
         *_calc_delta(bl_tot, bl_prev, brand, True)),
        ("RD Mới / Mất",       f"+{rd_new} / −{rd_lost}",
         f"▲ +{rd_new - rd_lost} ròng", get_color(brand, "success")),
    ]
    _mini_kpi_zone1(slide, MARGIN, _ZTY, _ZLW, _ZTH,
                    "🔗", "Chỉ số Backlink & Authority", kpis_z1, brand)

    # --- ZONE 2 (top-right): Backlinks nổi bật ---
    _add_rect(slide, _ZRX, _ZTY, _ZLW, _ZTH, "#F8F9FA")
    _zone_hdr(slide, _ZRX, _ZTY, _ZLW, "📥", "Backlinks nổi bật tháng này", primary, brand)

    cy = _ZTY + Inches(0.48)
    _add_text(slide, _ZRX + Inches(0.12), cy, _ZLW - Inches(0.2), Inches(0.24),
              "✅  Backlinks mới chất lượng cao",
              font_name=font_h, size=10, bold=True, color_hex=success)
    cy += Inches(0.26)
    for bl in ahrefs_data.get("new_notable_backlinks", [])[:3]:
        row_txt = f"  ↑  {bl['source']}  (DR {bl['dr']})  →  {bl['target']}"
        _add_text(slide, _ZRX + Inches(0.12), cy, _ZLW - Inches(0.2), Inches(0.28),
                  row_txt, font_name=font_b, size=10, color_hex=text_dark)
        cy += Inches(0.28)

    cy += Inches(0.1)
    _add_text(slide, _ZRX + Inches(0.12), cy, _ZLW - Inches(0.2), Inches(0.24),
              "⚠️  Backlinks mất trong tháng",
              font_name=font_h, size=10, bold=True, color_hex=warning)
    cy += Inches(0.26)
    for bl in ahrefs_data.get("lost_backlinks", [])[:2]:
        row_txt = f"  ↓  {bl['source']}  (DR {bl['dr']})  —  {bl['reason']}"
        _add_text(slide, _ZRX + Inches(0.12), cy, _ZLW - Inches(0.2), Inches(0.28),
                  row_txt, font_name=font_b, size=10, color_hex=text_light)
        cy += Inches(0.28)

    # --- ZONE 3 (bottom-left ✅): Điểm tốt ---
    _add_rect(slide, MARGIN, _ZBY, _ZLW, _ZBH, "#EFF8F1")
    _zone_hdr(slide, MARGIN, _ZBY, _ZLW, "✅", "Điểm tốt — Backlinks", success, brand)
    good_bullets = [
        f"DR tăng {dr_prev} → {dr_cur} (+3 điểm) — mức tăng DR tốt nhất trong 12 tháng",
        f"Referring domains tăng ròng +33: {rd_prev} → {rd_cur} (38 mới − 5 mất)",
        "5 backlink từ tên miền DR>60: vnexpress (78), cafebiz (72), brandsvietnam (65)",
        "74% dofollow — profile backlink tự nhiên, Google đánh giá tốt",
    ]
    by = _ZBY + Inches(0.5)
    for b in good_bullets:
        _add_text(slide, MARGIN + Inches(0.15), by, _ZLW - Inches(0.22), Inches(0.5),
                  f"•  {b}", font_name=font_b, size=10, color_hex=text_dark, wrap=True)
        by += Inches(0.5)

    # --- ZONE 4 (bottom-right ⚠️): Cần chú ý ---
    _add_rect(slide, _ZRX, _ZBY, _ZLW, _ZBH, "#FFFBF0")
    _zone_hdr(slide, _ZRX, _ZBY, _ZLW, "⚠️", "Cần chú ý — Backlinks", warning, brand)
    bad_bullets = [
        "5 referring domains mất — kiểm tra lý do, liên hệ webmaster nếu do lỗi kỹ thuật",
        "Đối thủ seo.net.vn (DR 51, 380 RDs) vẫn cách xa — cần duy trì momentum T5",
        "Ưu tiên T5: outreach 2-3 trang báo DR>50, tạo linkable asset mới",
        "Mục tiêu dài hạn: DR 50+ cuối năm 2026 để cạnh tranh top 3 ngành SEO",
    ]
    by2 = _ZBY + Inches(0.5)
    for b in bad_bullets:
        _add_text(slide, _ZRX + Inches(0.15), by2, _ZLW - Inches(0.22), Inches(0.5),
                  f"•  {b}", font_name=font_b, size=10, color_hex=text_dark, wrap=True)
        by2 += Inches(0.5)

    metrics_by_id = {m["id"]: m for m in comparison_data.get("metrics", [])}
    narr_raw  = metrics_by_id.get("domain_rating", {}).get(
        "interpretation",
        f"DR tăng {dr_prev} → {dr_cur} (+3 điểm). 33 referring domains mới ròng — momentum T4 rất tốt.")
    _narr(slide, brand, narr_raw[:135])
    _footer(slide, brand)


def create_slide_12_completed(prs, brand, actions_data, period_vn):
    """Slide 12: Việc đã hoàn thành trong tháng — bảng toàn trang."""
    slide = _blank_slide(prs)
    _slide_base(slide, brand, f"Việc đã hoàn thành — {period_vn}")

    success = get_color(brand, "success")

    tasks = actions_data.get("completed_this_month", [])
    headers     = ["#", "Hành động đã thực hiện", "Impact / Kết quả", "Owner"]
    col_widths  = [0.4, 6.0, 3.8, 2.133]
    rows        = []
    for i, t in enumerate(tasks, 1):
        impact = t.get("impact", "")
        if len(impact) > 85:
            impact = impact[:83] + "…"
        rows.append([str(i), t.get("task", ""), impact, t.get("owner", "")])

    table_y    = Inches(1.22)
    table_h    = H - Inches(1.22) - Inches(0.44)
    n_rows     = max(1, len(rows))
    table_h_in = table_h / 914400
    row_h      = max(0.40, (table_h_in - 0.38) / n_rows)

    _data_table(slide, MARGIN, table_y, W - 2 * MARGIN, table_h,
                headers, rows, brand,
                header_color=success, col_widths=col_widths, row_h_in=row_h)
    _footer(slide, brand)


def create_slide_13_plan(prs, brand, actions_data, next_period_vn):
    """Slide 13: Kế hoạch hành động tháng tới — bảng toàn trang."""
    slide = _blank_slide(prs)
    _slide_base(slide, brand, f"Kế hoạch hành động — {next_period_vn}")

    primary = get_color(brand, "primary")
    danger  = get_color(brand, "danger")
    warning = get_color(brand, "warning")
    t_light = get_color(brand, "text_light")

    tasks      = actions_data.get("in_progress_next_month", [])
    headers    = ["#", "Hành động kế hoạch", "Ưu tiên", "Owner", "Deadline"]
    col_widths = [0.4, 5.5, 1.2, 2.2, 3.033]
    p_colors   = {"P1": danger, "P2": warning, "P3": t_light}
    rows       = []
    for i, t in enumerate(tasks, 1):
        deadline = t.get("deadline", "")
        if deadline:
            parts    = deadline.split("-")
            deadline = f"{parts[2]}/{parts[1]}/{parts[0]}" if len(parts) == 3 else deadline
        rows.append([
            str(i), t.get("task", ""),
            t.get("priority", ""), t.get("owner", ""), deadline,
        ])

    def _pri_color(r, c, v):
        return p_colors.get(v) if c == 2 else None

    table_y    = Inches(1.22)
    table_h    = H - Inches(1.22) - Inches(0.44)
    n_rows     = max(1, len(rows))
    table_h_in = table_h / 914400
    row_h      = max(0.40, (table_h_in - 0.38) / n_rows)

    _data_table(slide, MARGIN, table_y, W - 2 * MARGIN, table_h,
                headers, rows, brand,
                header_color=primary, col_widths=col_widths, row_h_in=row_h,
                col_aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT,
                             PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER],
                cell_text_color_fn=_pri_color)
    _footer(slide, brand)


# ===========================================================================
# Build toàn bộ presentation
# ===========================================================================

def build_presentation(args) -> Path:
    brand_name = getattr(args, "brand", "default") or "default"
    brand      = load_brand(brand_name)

    if args.mock:
        data = load_mock_data()
    else:
        raise NotImplementedError("Live API chưa triển khai — dùng --mock")

    period_vn = _format_month_vn(args.month)
    today_str = date.today().strftime("%d/%m/%Y")

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("   📄 Slide  1: Trang bìa...")
    create_slide_1_cover(prs, brand, args.client, period_vn, args.domain, today_str)

    print("   📄 Slide  2: Mục lục...")
    create_slide_2_toc(prs, brand)

    print("   📄 Slide  3: Tóm tắt điều hành...")
    create_slide_3_exec_summary(prs, brand, data["comparison"])

    print("   📄 Slide  4: [Divider] I — Hiệu suất tổng quan...")
    create_slide_divider(prs, brand, "I", "Hiệu suất tổng quan",
                         "KPI chính  •  Xu hướng lưu lượng")

    print("   📄 Slide  5: KPI Summary Dashboard...")
    create_slide_5_kpi_dashboard(prs, brand, data["ga4"], data["gsc"])

    print("   📄 Slide  6: Xu hướng lưu lượng & chuyển đổi...")
    create_slide_6_traffic_trend(prs, brand, data["ga4"], data["gsc"], data["comparison"])

    print("   📄 Slide  7: [Divider] II — Phân tích chi tiết...")
    create_slide_divider(prs, brand, "II", "Phân tích chi tiết",
                         "GA4  •  GSC  •  Backlinks  •  4 góc nhìn")

    print("   📄 Slide  8: Phân tích GA4...")
    create_slide_8_ga4(prs, brand, data["ga4"], data["comparison"])

    print("   📄 Slide  9: Phân tích GSC...")
    create_slide_9_gsc(prs, brand, data["gsc"], data["comparison"])

    print("   📄 Slide 10: Backlinks & Domain Rating...")
    create_slide_10_backlinks(prs, brand, data["ahrefs"], data["comparison"])

    print("   📄 Slide 11: [Divider] III — Kế hoạch hành động...")
    create_slide_divider(prs, brand, "III", "Kế hoạch hành động",
                         "Việc đã hoàn thành  •  Kế hoạch tháng tới")

    print("   📄 Slide 12: Việc đã hoàn thành trong tháng...")
    create_slide_12_completed(prs, brand, data["actions"], period_vn)

    print("   📄 Slide 13: Kế hoạch hành động tháng tới...")
    next_period = _format_month_vn(
        f"{args.month[:4]}-{int(args.month[5:7]) + 1:02d}"
        if int(args.month[5:7]) < 12
        else f"{int(args.month[:4]) + 1}-01"
    )
    create_slide_13_plan(prs, brand, data["actions"], next_period)

    print("   📄 Slide 14: Khuyến nghị chiến lược dài hạn...")
    create_slide_14_longterm(prs, brand, data["actions"])

    print("   📄 Slide 15: Cảm ơn / Liên hệ...")
    create_slide_15_thankyou(prs, brand, args.domain)

    output_dir = Path(args.output_dir) if getattr(args, "output_dir", None) else OUTPUTS_DIR
    output_dir.mkdir(parents=True, exist_ok=True)
    client_safe = re.sub(r"[^a-zA-Z0-9_]", "_", args.client)
    output_path = output_dir / f"{client_safe}_SEO_Report_{args.month}.pptx"
    prs.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def parse_args():
    p = argparse.ArgumentParser(description="Tạo monthly SEO report .pptx 15 slides")
    p.add_argument("--client",     required=True,   help="Tên client")
    p.add_argument("--domain",     required=True,   help="Domain (vd: viettelstore.vn)")
    p.add_argument("--month",      required=True,   help="Tháng YYYY-MM (vd: 2026-04)")
    p.add_argument("--brand",      default="default", help="Brand profile (mặc định: default)")
    p.add_argument("--mock",       action="store_true", help="Dùng mock data")
    p.add_argument("--output-dir", dest="output_dir",   help="Thư mục output")
    return p.parse_args()


def main():
    args = parse_args()
    print(f"\n{'='*60}")
    print(f"  SEONGON Monthly SEO Report Generator")
    print(f"{'='*60}")
    print(f"  Client  : {args.client}")
    print(f"  Domain  : {args.domain}")
    print(f"  Tháng   : {args.month}")
    print(f"  Brand   : {args.brand}")
    print(f"  Mode    : {'MOCK DATA' if args.mock else 'LIVE API'}")
    print(f"{'='*60}\n")

    output_path = build_presentation(args)
    size_kb     = output_path.stat().st_size // 1024

    print(f"\n{'='*60}")
    print(f"  ✅ Hoàn thành!")
    print(f"  File   : {output_path}")
    print(f"  Size   : {size_kb} KB")
    print(f"  Slides : 15")
    print(f"{'='*60}\n")


if __name__ == "__main__":
    main()
